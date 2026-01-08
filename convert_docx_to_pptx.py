import sys
import os
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def convert_docx_to_pptx(docx_file, pptx_file):
    if not os.path.isfile(docx_file):
        raise FileNotFoundError(f"Input file '{docx_file}' not found.")
    
    if not docx_file.lower().endswith('.docx'):
        raise ValueError("Error: Input file must be a .docx file.")

    # Ensure the output directory exists
    # output_dir = os.path.dirname(pptx_file)
    # if not os.path.exists(output_dir):
    #     os.makedirs(output_dir)

    doc = Document(docx_file)
    presentation = Presentation()

    for para in doc.paragraphs:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        left = Inches(1)
        top = Inches(1)
        width = Inches(8.5)
        height = Inches(5.5)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame

        p = text_frame.add_paragraph()
        p.text = para.text
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)

        # Remove empty paragraphs if any
        if len(text_frame.paragraphs) > 1:
            for paragraph in text_frame.paragraphs[:-1]:
                if not paragraph.text.strip():
                    p._element.getparent().remove(p._element)

    presentation.save(pptx_file)

def rename_and_convert(input_file, output_file):
    # Rename input file to have .docx extension if it doesn't already
    if not input_file.lower().endswith('.docx'):
        temp_docx_file = f"{input_file}.docx"
        os.rename(input_file, temp_docx_file)
    else:
        temp_docx_file = input_file

    try:
        # Perform the conversion
        convert_docx_to_pptx(temp_docx_file, output_file)
    finally:
        # Revert the file name back to the original
        if not input_file.lower().endswith('.docx'):
            os.rename(temp_docx_file, input_file)

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python convert_docx_to_pptx.py <input_docx_file> <output_pptx_file>")
        sys.exit(1)

    docx_file = sys.argv[1]
    pptx_file = sys.argv[2]

    try:
        rename_and_convert(docx_file, pptx_file)
        print(f"Conversion complete. PPTX saved to {pptx_file}")
    except Exception as e:
        print(f"Error during conversion: {e}")
